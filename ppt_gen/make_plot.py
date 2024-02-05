# %%
# python script which loads png images from a directory and creates a powerpoint slide for every two images. In the powerpoint slide, the two images are placed side by side.


import os
import pptx
from pathlib import Path

image_dir = "/mnt/dl-41/data/leeuwenmcv/general/l3harris/yolov8l_dist_est_03022024"
images = list(Path(image_dir).rglob("*wdist.png"))
images.sort()

prs = pptx.Presentation()
for i in range(0, len(images), 2):
	slide = prs.slides.add_slide(prs.slide_layouts[5])
	slide_width = prs.slide_width
	slide_height = prs.slide_height

	# Calculate the maximum width and height for the images
	max_width = int(slide_width * 0.45)
	max_height = int(slide_height * 0.8)

	# Add the left image
	left = slide.shapes.add_picture(str(images[i]), 0, 0, max_width, max_height)

	if i + 1 < len(images):
		# Add the right image
		right = slide.shapes.add_picture(
			str(images[i + 1]), max_width, 0, max_width, max_height
		)

prs.save("output.pptx")
